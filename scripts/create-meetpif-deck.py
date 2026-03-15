#!/usr/bin/env python3
"""
Create an example slide deck using the meetpif.com design system.
Demonstrates: title, section divider, content, data, comparison, quote, CTA slides.
"""

import sys
sys.path.insert(0, '/root/projects/duvo/deck')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx_helpers import set_slide_bg, add_rect, add_rounded_rect, add_circle, \
    add_text_box, add_multiline_text

# ─── Brand tokens ───────────────────────────────────────────────────────────
BLACK = RGBColor(0x00, 0x00, 0x00)
SURFACE = RGBColor(0x11, 0x11, 0x11)
SURFACE_RAISED = RGBColor(0x1A, 0x1A, 0x1A)
LIME = RGBColor(0xD8, 0xFF, 0x66)
LIME_MUTED = RGBColor(0x2A, 0x33, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SECONDARY_TEXT = RGBColor(0x99, 0x99, 0x99)
TERTIARY_TEXT = RGBColor(0x55, 0x55, 0x55)
BORDER = RGBColor(0x22, 0x22, 0x22)
SUCCESS = RGBColor(0x66, 0xFF, 0xB2)
ERROR = RGBColor(0xFF, 0x6B, 0x6B)
WARNING = RGBColor(0xFF, 0xB8, 0x4D)
INFO = RGBColor(0x66, 0xD9, 0xFF)

# Fonts
FONT_HEADLINE = "Inter"  # python-pptx can't guarantee Inter Tight in all envs, Inter is close
FONT_BODY = "Inter"
FONT_DISPLAY = "Georgia"  # Fallback for Instrument Serif

# Slide dimensions: 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# Helpers imported from pptx_helpers module.


# ─── Build the deck ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # Blank

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1: Title — Hero style (oversized text + central logo, inspired by meetpif.com)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

# Oversized "PIF LABORMAN —" text spanning full width, vertically centered
# The hero uses clamp(100px, 18vw, 280px) — at 1920px that's ~280px ≈ 210pt
# We use a wide text box that extends beyond slide edges for the bleed effect
marquee_text = "PIF LABORMAN —"
# Text box wider than slide, centered, to simulate edge-bleed
overflow_w = Inches(18)  # wider than 13.333" slide
overflow_left = Inches((13.333 - 18) / 2)  # negative offset to center
add_text_box(slide, overflow_left, Inches(2.5), overflow_w, Inches(3),
             marquee_text, FONT_HEADLINE, 160, WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, letter_spacing=-3)

# Central lime green rounded rectangle (the Pif avatar/logo shape)
# On the hero it's ~65vh = about 4.8" on a 7.5" slide
logo_size = Inches(4.2)
logo_x = Inches((13.333 - 4.2) / 2)
logo_y = Inches((7.5 - 4.2) / 2) - Inches(0.15)  # slightly above center
avatar = add_rounded_rect(slide, logo_x, logo_y, logo_size, logo_size, LIME)
# Adjust corner rounding for the avatar shape
avatar.adjustments[0] = 0.12  # ~12% radius, creates the app-icon look

# Two "eyes" on the avatar (dark circles)
eye_size = Inches(0.55)
eye_y = logo_y + Inches(1.8)
eye_left_x = logo_x + Inches(1.35)
eye_right_x = logo_x + Inches(2.3)
add_circle(slide, eye_left_x, eye_y, eye_size, BLACK)
add_circle(slide, eye_right_x, eye_y, eye_size, BLACK)

# "DESIGN SYSTEM" subtitle — below the marquee band, left-aligned
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(5), Inches(0.5),
             "DESIGN SYSTEM", FONT_HEADLINE, 14, LIME, bold=True,
             letter_spacing=3)

# Bottom-left: operational text
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(5), Inches(0.3),
             "MARCH 2026", FONT_HEADLINE, 11, TERTIARY_TEXT,
             letter_spacing=2)

# Bottom-right: tagline (mimicking the hero's asymmetric bottom metadata)
add_text_box(slide, Inches(7), Inches(6.2), Inches(5.5), Inches(0.4),
             "Brand Guidelines & Tokens", FONT_BODY, 20, SECONDARY_TEXT,
             alignment=PP_ALIGN.RIGHT)
add_text_box(slide, Inches(7), Inches(6.65), Inches(5.5), Inches(0.35),
             "Colors · Typography · Components · Spacing",
             FONT_BODY, 14, TERTIARY_TEXT, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(5), Inches(0.6),
             "AGENDA", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

# Divider line
add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

agenda_items = [
    ("01", "Design Identity"),
    ("02", "Color System"),
    ("03", "Typography Scale"),
    ("04", "Component Library"),
    ("05", "Layout & Spacing"),
    ("06", "Getting Started"),
]

for i, (num, label) in enumerate(agenda_items):
    y = Inches(2.0 + i * 0.85)
    # Number
    add_text_box(slide, Inches(1), y, Inches(0.8), Inches(0.5),
                 num, FONT_HEADLINE, 28, LIME, bold=True)
    # Label
    add_text_box(slide, Inches(2), y, Inches(8), Inches(0.5),
                 label, FONT_BODY, 20, WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3: Section Divider — Design Identity
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, LIME)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "01", FONT_HEADLINE, 80, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=4)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
             "DESIGN IDENTITY", FONT_HEADLINE, 44, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4: Design Identity — Principles
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "DESIGN PRINCIPLES", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

principles = [
    ("DARK-FIRST", "The primary experience is dark mode. Light is the alternate."),
    ("UPPERCASE COMMANDS", "All headings, buttons, labels — ALL CAPS with wide tracking."),
    ("LIME FOR ACTION", "Primary buttons, active states, focus rings use accent green."),
    ("SEMANTIC ONLY", "Success, error, warning, info — for feedback, never decoration."),
    ("SHARP CORNERS", "4px buttons, 6px inputs, 10px cards. No pills except badges."),
]

for i, (title, desc) in enumerate(principles):
    y = Inches(2.0 + i * 1.0)
    # Lime dot
    add_circle(slide, Inches(1), y + Inches(0.08), Inches(0.18), LIME)
    # Title
    add_text_box(slide, Inches(1.5), y, Inches(4), Inches(0.4),
                 title, FONT_HEADLINE, 16, WHITE, bold=True, letter_spacing=2)
    # Description
    add_text_box(slide, Inches(1.5), y + Inches(0.38), Inches(10), Inches(0.4),
                 desc, FONT_BODY, 14, SECONDARY_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5: Section Divider — Color System
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, LIME)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "02", FONT_HEADLINE, 80, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=4)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
             "COLOR SYSTEM", FONT_HEADLINE, 44, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6: Color Palette
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "COLOR PALETTE", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

# Color swatches
colors = [
    (LIME, "#D8FF66", "ACCENT PRIMARY", "Primary actions, focus, active states"),
    (WHITE, "#FFFFFF", "TEXT PRIMARY", "Headings, body text, labels"),
    (SECONDARY_TEXT, "#999999", "TEXT SECONDARY", "Descriptions, placeholders"),
    (SURFACE, "#111111", "SURFACE", "Cards, panels, raised elements"),
    (BLACK, "#000000", "PAGE BG", "Default page background"),
]

for i, (color, hex_val, name, usage) in enumerate(colors):
    x = Inches(1 + i * 2.3)
    y = Inches(2.2)

    # Color swatch
    swatch = add_rounded_rect(slide, x, y, Inches(1.8), Inches(1.8), color,
                               BORDER if color == BLACK else None)

    # Hex value
    text_color = BLACK if color in (LIME, WHITE) else WHITE
    add_text_box(slide, x, y + Inches(2.0), Inches(1.8), Inches(0.35),
                 hex_val, FONT_BODY, 13, SECONDARY_TEXT, alignment=PP_ALIGN.CENTER)

    # Name
    add_text_box(slide, x, y + Inches(2.4), Inches(1.8), Inches(0.35),
                 name, FONT_HEADLINE, 11, WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, uppercase=True, letter_spacing=1.5)

    # Usage
    add_text_box(slide, x, y + Inches(2.8), Inches(1.8), Inches(0.5),
                 usage, FONT_BODY, 10, TERTIARY_TEXT, alignment=PP_ALIGN.CENTER)

# Semantic colors row
add_text_box(slide, Inches(1), Inches(5.6), Inches(4), Inches(0.4),
             "SEMANTIC COLORS", FONT_HEADLINE, 14, WHITE, bold=True, letter_spacing=2)

semantic = [
    (SUCCESS, "SUCCESS", "#66FFB2"),
    (ERROR, "ERROR", "#FF6B6B"),
    (WARNING, "WARNING", "#FFB84D"),
    (INFO, "INFO", "#66D9FF"),
]

for i, (color, name, hex_val) in enumerate(semantic):
    x = Inches(1 + i * 2.3)
    y = Inches(6.1)
    add_circle(slide, x, y, Inches(0.4), color)
    add_text_box(slide, x + Inches(0.55), y + Inches(0.02), Inches(1.5), Inches(0.2),
                 name, FONT_HEADLINE, 11, WHITE, bold=True, letter_spacing=1.5)
    add_text_box(slide, x + Inches(0.55), y + Inches(0.22), Inches(1.5), Inches(0.2),
                 hex_val, FONT_BODY, 10, TERTIARY_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7: Section Divider — Typography
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, LIME)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "03", FONT_HEADLINE, 80, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=4)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
             "TYPOGRAPHY", FONT_HEADLINE, 44, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8: Typography Scale
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "TYPE SCALE", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

# Left column: The "loud" font
add_text_box(slide, Inches(1), Inches(2.0), Inches(5), Inches(0.4),
             'INTER TIGHT — "THE LOUD FONT"', FONT_HEADLINE, 12, LIME, bold=True, letter_spacing=2)

type_samples_loud = [
    ("DISPLAY LARGE", 36, True),
    ("HEADING ONE", 24, True),
    ("HEADING TWO", 20, True),
    ("BUTTON LABEL", 12, True),
]

y_pos = 2.6
for text, size, bold in type_samples_loud:
    add_text_box(slide, Inches(1), Inches(y_pos), Inches(5.5), Inches(0.8),
                 text, FONT_HEADLINE, size, WHITE, bold=bold, letter_spacing=2)
    add_text_box(slide, Inches(1), Inches(y_pos + 0.5), Inches(5.5), Inches(0.3),
                 f"{size}px / Bold / ALL CAPS / Wide tracking", FONT_BODY, 10, TERTIARY_TEXT)
    y_pos += 0.95

# Right column: The "quiet" font
add_text_box(slide, Inches(7), Inches(2.0), Inches(5), Inches(0.4),
             'INTER — "THE QUIET FONT"', FONT_HEADLINE, 12, LIME, bold=True, letter_spacing=2)

type_samples_quiet = [
    ("Body text at 16px for prominent content", 16, False),
    ("Default body text at 14px for readable content", 14, False),
    ("Small text at 13px for inputs and descriptions", 13, False),
    ("Caption text at 12px for timestamps", 12, False),
]

y_pos = 2.6
for text, size, bold in type_samples_quiet:
    add_text_box(slide, Inches(7), Inches(y_pos), Inches(5.5), Inches(0.8),
                 text, FONT_BODY, size, WHITE, bold=bold)
    add_text_box(slide, Inches(7), Inches(y_pos + 0.4), Inches(5.5), Inches(0.3),
                 f"{size}px / Regular / Sentence case", FONT_BODY, 10, TERTIARY_TEXT)
    y_pos += 0.75

# Display font note
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Instrument Serif for decorative/editorial display text only. JetBrains Mono for code blocks.",
             FONT_BODY, 12, TERTIARY_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9: Section Divider — Components
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, LIME)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "04", FONT_HEADLINE, 80, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=4)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
             "COMPONENTS", FONT_HEADLINE, 44, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10: Component Overview — 141 components
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "141 COMPONENTS", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

# Component category cards
categories = [
    ("BUTTONS", "28"),
    ("FORMS", "21"),
    ("TYPOGRAPHY", "16"),
    ("NAVIGATION", "12"),
    ("BADGES", "9"),
    ("FEEDBACK", "9"),
    ("CARDS", "7"),
    ("INDICATORS", "6"),
    ("DATA", "5"),
    ("MENUS", "4"),
    ("LAYOUT", "4"),
    ("SKELETONS", "4"),
    ("OVERLAYS", "3"),
    ("AVATARS", "3"),
    ("FILTERS", "3"),
    ("COMMENTS", "2"),
]

cols = 4
for i, (name, count) in enumerate(categories):
    col = i % cols
    row = i // cols
    x = Inches(1 + col * 2.9)
    y = Inches(2.0 + row * 1.3)

    # Card bg
    add_rounded_rect(slide, x, y, Inches(2.5), Inches(1.0), SURFACE, BORDER)

    # Count
    add_text_box(slide, x + Inches(0.2), y + Inches(0.12), Inches(1), Inches(0.5),
                 count, FONT_HEADLINE, 28, LIME, bold=True)

    # Name
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(2), Inches(0.3),
                 name, FONT_HEADLINE, 11, SECONDARY_TEXT, bold=True, letter_spacing=1.5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11: Button Examples
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "BUTTONS", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

add_text_box(slide, Inches(1), Inches(1.8), Inches(10), Inches(0.4),
             "Inter Tight 12px / 500 weight / ALL CAPS / 1.2px tracking / 4px border radius",
             FONT_BODY, 12, TERTIARY_TEXT)

# Primary button
btn = add_rounded_rect(slide, Inches(1), Inches(2.5), Inches(2.4), Inches(0.55), LIME)
add_text_box(slide, Inches(1), Inches(2.55), Inches(2.4), Inches(0.45),
             "PRIMARY ACTION", FONT_HEADLINE, 12, BLACK, bold=True,
             alignment=PP_ALIGN.CENTER, letter_spacing=1.5)

# Secondary button
btn = add_rounded_rect(slide, Inches(3.8), Inches(2.5), Inches(2.4), Inches(0.55), SURFACE, BORDER)
add_text_box(slide, Inches(3.8), Inches(2.55), Inches(2.4), Inches(0.45),
             "SECONDARY", FONT_HEADLINE, 12, WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, letter_spacing=1.5)

# Ghost button
add_text_box(slide, Inches(6.6), Inches(2.55), Inches(2.4), Inches(0.45),
             "GHOST BUTTON", FONT_HEADLINE, 12, WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, letter_spacing=1.5)

# Danger button
btn = add_rounded_rect(slide, Inches(9.4), Inches(2.5), Inches(2.4), Inches(0.55), ERROR)
add_text_box(slide, Inches(9.4), Inches(2.55), Inches(2.4), Inches(0.45),
             "DANGER", FONT_HEADLINE, 12, WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, letter_spacing=1.5)

# Button specs text
specs = [
    "Primary: lime green fill (#D8FF66), black text, no border",
    "Secondary: surface fill (#111111), white text, 1px border",
    "Ghost: transparent, white text, no border",
    "Danger: error red fill (#FF6B6B), white text",
    "",
    "Padding: 12px 24px (default) / 8px 16px (small)",
    "States: hover, disabled (opacity 0.5), loading (spinner + label change)",
    "Icon buttons: 32x32 ghost or subtle, 6px radius",
]

add_multiline_text(slide, Inches(1), Inches(3.5), Inches(10), Inches(3.5),
                   specs, FONT_BODY, 14, SECONDARY_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12: Section Divider — Layout
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, LIME)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "05", FONT_HEADLINE, 80, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=4)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
             "LAYOUT & SPACING", FONT_HEADLINE, 44, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13: Spacing Scale
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "SPACING SCALE", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

spacing_tokens = [
    ("XS", "4px", 0.06),
    ("SM", "8px", 0.12),
    ("MD", "12px", 0.18),
    ("LG", "16px", 0.24),
    ("XL", "24px", 0.36),
    ("2XL", "32px", 0.48),
    ("3XL", "48px", 0.72),
    ("4XL", "64px", 0.96),
]

for i, (name, value, bar_w) in enumerate(spacing_tokens):
    y = Inches(2.0 + i * 0.6)
    # Label
    add_text_box(slide, Inches(1), y, Inches(0.8), Inches(0.35),
                 name, FONT_HEADLINE, 12, WHITE, bold=True, letter_spacing=1.5)
    # Value
    add_text_box(slide, Inches(1.8), y, Inches(0.8), Inches(0.35),
                 value, FONT_BODY, 12, SECONDARY_TEXT)
    # Visual bar
    add_rounded_rect(slide, Inches(3), y + Inches(0.05), Inches(bar_w), Inches(0.25), LIME)

# Border radius section on right
add_text_box(slide, Inches(7), Inches(2.0), Inches(5), Inches(0.4),
             "BORDER RADIUS", FONT_HEADLINE, 14, LIME, bold=True, letter_spacing=2)

radii = [
    ("Buttons", "4px", "--radius-sm"),
    ("Inputs", "6px", "--radius-md"),
    ("Dropdowns", "8px", "--radius-lg"),
    ("Cards", "10px", "custom"),
    ("Modals", "14px", "custom"),
    ("Badges", "pill", "--radius-pill"),
]

for i, (component, value, token) in enumerate(radii):
    y = Inches(2.6 + i * 0.55)
    add_text_box(slide, Inches(7), y, Inches(2), Inches(0.3),
                 component, FONT_BODY, 13, WHITE)
    add_text_box(slide, Inches(9.2), y, Inches(1), Inches(0.3),
                 value, FONT_BODY, 13, LIME, bold=True)
    add_text_box(slide, Inches(10.3), y, Inches(2), Inches(0.3),
                 token, FONT_BODY, 11, TERTIARY_TEXT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14: Section Divider — Getting Started
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, LIME)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
             "06", FONT_HEADLINE, 80, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=4)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(1),
             "GETTING STARTED", FONT_HEADLINE, 44, BLACK, bold=True,
             alignment=PP_ALIGN.LEFT, letter_spacing=5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15: Quick Start Rules
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

add_text_box(slide, Inches(1), Inches(0.8), Inches(6), Inches(0.6),
             "RULES FOR BUILDERS", FONT_HEADLINE, 24, LIME, bold=True, letter_spacing=3)

add_rect(slide, Inches(1), Inches(1.5), Inches(2), Inches(0.02), LIME)

rules = [
    "Use CSS variables from tokens.css — never hardcode values",
    "Check components.md before building any component",
    "Headlines and buttons are ALWAYS uppercase with tracking",
    "Inter Tight for headlines/buttons, Inter for body",
    "Primary buttons use lime green, not black or white",
    "Buttons use 4px radius — only badges are pill-shaped",
    "Dark theme is default — :root styles are dark",
    "Badge text is uppercase with colored dot indicators",
    "Modal sections use 1px dividers, not just spacing",
    "Tab active indicator is 2px lime green bar below label",
]

for i, rule in enumerate(rules):
    y = Inches(1.9 + i * 0.5)
    # Number
    num = f"{i+1:02d}"
    add_text_box(slide, Inches(1), y, Inches(0.6), Inches(0.4),
                 num, FONT_HEADLINE, 14, LIME, bold=True)
    # Rule text
    add_text_box(slide, Inches(1.6), y, Inches(10), Inches(0.4),
                 rule, FONT_BODY, 15, WHITE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16: Closing / CTA — Hero style bookend
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BLACK)

# Same oversized text treatment as title slide
add_text_box(slide, overflow_left, Inches(2.5), overflow_w, Inches(3),
             "MEETPIF —", FONT_HEADLINE, 160, WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, letter_spacing=-3)

# Central lime avatar (smaller than title — visual variety)
logo_close_size = Inches(3.2)
logo_close_x = Inches((13.333 - 3.2) / 2)
logo_close_y = Inches((7.5 - 3.2) / 2) - Inches(0.15)
avatar_close = add_rounded_rect(slide, logo_close_x, logo_close_y,
                                 logo_close_size, logo_close_size, LIME)
avatar_close.adjustments[0] = 0.12
# Eyes
eye_sm = Inches(0.42)
eye_close_y = logo_close_y + Inches(1.35)
add_circle(slide, logo_close_x + Inches(1.0), eye_close_y, eye_sm, BLACK)
add_circle(slide, logo_close_x + Inches(1.78), eye_close_y, eye_sm, BLACK)

# Bottom-left: file references
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(6), Inches(0.3),
             "tokens.css  ·  typography.md  ·  components.md  ·  spacing.md",
             FONT_BODY, 12, TERTIARY_TEXT)

# Bottom-right: URL + tagline
add_text_box(slide, Inches(7), Inches(6.0), Inches(5.5), Inches(0.4),
             "meetpif.com", FONT_HEADLINE, 18, LIME, bold=True,
             alignment=PP_ALIGN.RIGHT, letter_spacing=2)
add_text_box(slide, Inches(7), Inches(6.45), Inches(5.5), Inches(0.35),
             "Dark-first. Uppercase. Lime green.",
             FONT_BODY, 14, SECONDARY_TEXT, alignment=PP_ALIGN.RIGHT)


# ─── Save ───────────────────────────────────────────────────────────────────
output_path = "/root/data/meetpif-design-system-deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
