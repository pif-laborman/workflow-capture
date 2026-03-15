#!/usr/bin/env python3
"""
Duvo.ai brand presentation deck — generated from Figma design tokens.

Visual language extracted from Figma file (duvo-website, Pavol copy):
- Warm yellow radial glow backgrounds (illumination / energy)
- Yellow-to-mint gradient transitions
- Clean sans-serif typography (Inter family)
- White/light content cards with rounded corners on gradient backgrounds
- Duvo play-triangle logo with yellow accent dots
- Black square footer element, duvo.ai logo bottom-left
- Custom isometric illustrations (placeholder in programmatic version)
- Product badges: "Duvo Clarity", "Duvo Case Queue"
"""

import sys
sys.path.insert(0, '/root/projects/duvo/deck')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx_helpers import set_slide_bg, add_rect, add_rounded_rect, add_circle, \
    add_text_box, add_pill
import math

# ─── Duvo Brand Tokens ──────────────────────────────────────────────────────

# Core palette (from Figma branding + styleguide)
DUVO_YELLOW = RGBColor(0xFF, 0xC8, 0x00)       # Primary brand yellow
DUVO_YELLOW_LIGHT = RGBColor(0xFF, 0xE0, 0x4D)  # Lighter yellow for glows
DUVO_YELLOW_PALE = RGBColor(0xFF, 0xF3, 0xCC)   # Very pale yellow for bg wash
DUVO_YELLOW_BG = RGBColor(0xFE, 0xF9, 0xE8)     # Warm cream background
DUVO_MINT = RGBColor(0xE8, 0xF5, 0xF0)          # Mint for card backgrounds
DUVO_MINT_ACCENT = RGBColor(0xC8, 0xE6, 0xDB)   # Stronger mint for borders
DUVO_TEAL = RGBColor(0x1A, 0x8A, 0x7A)          # Teal for icons

# Backgrounds
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WARM = RGBColor(0xFD, 0xFA, 0xF2)            # Warm white base
BG_SURFACE = RGBColor(0xF7, 0xF5, 0xF0)         # Light surface
BG_DARK = RGBColor(0x1A, 0x1A, 0x1A)            # Dark sections
BG_BLACK = RGBColor(0x00, 0x00, 0x00)

# Text
TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SECONDARY = RGBColor(0x55, 0x55, 0x55)
TEXT_TERTIARY = RGBColor(0x88, 0x88, 0x88)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_YELLOW = RGBColor(0xCC, 0xA0, 0x00)         # Yellow-on-white for accents

# Semantic
SUCCESS_GREEN = RGBColor(0x03, 0xA9, 0x7E)
PRICE_GREEN = RGBColor(0x00, 0xC4, 0x8C)        # For pricing comparison

# Borders
BORDER_LIGHT = RGBColor(0xE8, 0xE5, 0xE0)
BORDER_CARD = RGBColor(0xEE, 0xEB, 0xE5)

# Fonts (from Figma styleguide)
FONT_HEADLINE = "Inter"
FONT_BODY = "Inter"

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ────────────────────────────────────────────────────────────────
# Core helpers (set_slide_bg, add_rect, add_rounded_rect, add_circle,
# add_text_box, add_pill) imported from pptx_helpers.
# add_multiline kept as alias with brand-specific default spacing.
from pptx_helpers import add_multiline_text as _add_multiline_text

def add_multiline(slide, left, top, width, height, lines, font_name=FONT_BODY,
                  font_size=14, color=TEXT_PRIMARY, alignment=PP_ALIGN.LEFT,
                  line_spacing=None, bold=False):
    return _add_multiline_text(slide, left, top, width, height, lines,
                               font_name=font_name, font_size=font_size,
                               color=color, alignment=alignment,
                               line_spacing=line_spacing or font_size * 0.4,
                               bold=bold)


def yellow_glow_bg(slide):
    """Simulate the Duvo warm yellow radial glow background.
    Warm cream base with yellow glow center-left and mint fade on right.
    """
    set_slide_bg(slide, BG_WARM)
    # Yellow radial glow effect (center-left, simulated with overlapping rects)
    add_rect(slide, Inches(2), Inches(1), Inches(5), Inches(5),
             DUVO_YELLOW_PALE)
    add_rect(slide, Inches(3), Inches(1.5), Inches(3.5), Inches(4),
             RGBColor(0xFF, 0xF0, 0xB0))
    add_rect(slide, Inches(3.5), Inches(2), Inches(2.5), Inches(3),
             RGBColor(0xFF, 0xE8, 0x80))
    # Mint fade on right edge
    add_rect(slide, Inches(9.5), Inches(0), Inches(4), SLIDE_H,
             RGBColor(0xF2, 0xF9, 0xF5))
    add_rect(slide, Inches(10.5), Inches(1), Inches(3), Inches(5),
             DUVO_MINT)


def content_card_bg(slide):
    """Light gradient background for content slides — subtle yellow-to-mint."""
    set_slide_bg(slide, BG_WHITE)
    # Subtle warm wash on left
    add_rect(slide, Inches(0), Inches(0), Inches(6), SLIDE_H,
             RGBColor(0xFE, 0xFC, 0xF5))
    # Subtle mint on right
    add_rect(slide, Inches(8), Inches(0), Inches(5.5), SLIDE_H,
             RGBColor(0xF5, 0xFB, 0xF8))


def duvo_footer(slide, dark=False):
    """Standard Duvo footer: logo text bottom-left, black square bottom-right."""
    txt_color = TEXT_WHITE if dark else TEXT_PRIMARY
    # Logo text
    add_text_box(slide, Inches(0.6), Inches(6.85), Inches(2), Inches(0.4),
                 "▶  duvo.ai", FONT_HEADLINE, 14, txt_color, bold=True)
    # Black square (brand element)
    sq_color = TEXT_WHITE if dark else BG_BLACK
    add_rect(slide, Inches(12.3), Inches(6.75), Inches(0.45), Inches(0.45), sq_color)


def section_badge(slide, left, top, text, inverted=False):
    """Small pill badge for product sections (e.g. 'Duvo Clarity')."""
    bg = DUVO_YELLOW if not inverted else RGBColor(0xFF, 0xE0, 0x4D)
    txt = TEXT_PRIMARY
    add_pill(slide, left, top, Inches(1.8), Inches(0.35), bg)
    # Play icon placeholder
    add_text_box(slide, left + Inches(0.1), top + Inches(0.02),
                 Inches(0.3), Inches(0.3), "▶", FONT_BODY, 10, txt)
    add_text_box(slide, left + Inches(0.35), top + Inches(0.04),
                 Inches(1.4), Inches(0.28), text, FONT_BODY, 11, txt, bold=True)


# ─── Build the deck ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER — Bold headline, yellow glow bg, golden swoosh placeholder
# Matches: retail-slide-01.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
yellow_glow_bg(slide)

# Logo top-left
add_text_box(slide, Inches(0.6), Inches(0.5), Inches(2.5), Inches(0.4),
             "▶  duvo.ai", FONT_HEADLINE, 18, TEXT_PRIMARY, bold=True)

# Golden swoosh placeholder (right side abstract curves)
# In Figma this is a vector illustration — we simulate with yellow arcs
for i in range(3):
    offset = i * 0.3
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(9 + offset), Inches(1.5 + offset),
        Inches(4.5 - offset), Inches(5 - offset * 2)
    )
    shape.fill.solid()
    yellow_val = 0xC8 + i * 20
    shape.fill.fore_color.rgb = RGBColor(0xFF, min(yellow_val, 0xFF), 0x00)
    shape.line.fill.background()
    shape.rotation = -15 + i * 10

# Main headline — big, bold, black
add_text_box(slide, Inches(0.6), Inches(3.0), Inches(9), Inches(2.5),
             "Duvo removes the manual\nbusywork from retail ops.",
             FONT_HEADLINE, 48, TEXT_PRIMARY, bold=True)

# Subtitle
add_text_box(slide, Inches(0.6), Inches(5.4), Inches(8), Inches(0.6),
             "End-to-end execution across your existing tools for a fraction of\nthe manual cost.",
             FONT_BODY, 16, TEXT_SECONDARY)

# Date
add_text_box(slide, Inches(0.6), Inches(6.3), Inches(3), Inches(0.3),
             "2026 MARCH", FONT_BODY, 12, TEXT_TERTIARY)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM — "More than half your team's time is lost"
# Matches: retail-slide-02.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

# Headline
add_text_box(slide, Inches(0.6), Inches(0.5), Inches(10), Inches(1.2),
             "More than half of your team's time is\nlost to cross-system busywork.",
             FONT_HEADLINE, 36, TEXT_PRIMARY, bold=True)

# Left callout
add_circle(slide, Inches(0.8), Inches(2.2), Inches(0.5), DUVO_MINT)
add_text_box(slide, Inches(0.6), Inches(2.9), Inches(4.5), Inches(1.5),
             "Smart people hired to\ngrow the business are\nstuck bridging its\ndisconnected systems.",
             FONT_HEADLINE, 22, TEXT_PRIMARY, bold=True)

# Right content card (mint background)
add_rounded_rect(slide, Inches(5.5), Inches(2.0), Inches(7), Inches(4.5),
                 DUVO_MINT, DUVO_MINT_ACCENT)

add_text_box(slide, Inches(5.9), Inches(2.3), Inches(6.2), Inches(0.6),
             "In practice, this results in often repetitive, manual work:",
             FONT_BODY, 14, TEXT_SECONDARY)

manual_tasks = [
    ("…moving information across tools\n(ERP, portals, email, BI…)", Inches(3.2)),
    ("…reconciling spreadsheets\nwith price lists and inventory files", Inches(4.3)),
    ("…translating insights into actions\nthat still need manual execution", Inches(5.3)),
]
for text, y in manual_tasks:
    add_circle(slide, Inches(9.5), y, Inches(0.3), DUVO_TEAL)
    add_text_box(slide, Inches(6.0), y, Inches(5), Inches(0.8),
                 text, FONT_BODY, 13, TEXT_SECONDARY)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM VISUAL — "Signals everywhere. Exceptions piling up."
# Matches: retail-slide-03.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(1.2),
             "Signals everywhere. Exceptions piling up.\nYour team is juggling it all.",
             FONT_HEADLINE, 36, TEXT_PRIMARY, bold=True)

# Two column labels
add_text_box(slide, Inches(1), Inches(2.0), Inches(5), Inches(0.6),
             "Signals come from everywhere\nDashboards, reports, alerts, exceptions",
             FONT_BODY, 13, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(7), Inches(2.0), Inches(5), Inches(0.6),
             "Execution happens across disconnected tools\nERPs, CRMs, calls, emails, spreadsheets",
             FONT_BODY, 13, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Illustration placeholder (isometric character juggling systems)
add_rounded_rect(slide, Inches(2), Inches(2.8), Inches(9), Inches(3.5),
                 BG_SURFACE, BORDER_LIGHT)
add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(1),
             "[Illustration: person juggling disconnected systems —\nSAP, Oracle, emails, spreadsheets, ERPs]",
             FONT_BODY, 14, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

# Bottom label
add_text_box(slide, Inches(3), Inches(6.1), Inches(7), Inches(0.4),
             "Your team is the glue\nCopying, chasing, re-keying – manually.",
             FONT_BODY, 13, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4: COST — "The real cost shows up in three places"
# Matches: retail-slide-04.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.5), Inches(11), Inches(1.2),
             "The real cost of this bottleneck\nshows up in three places.",
             FONT_HEADLINE, 36, TEXT_PRIMARY, bold=True)

costs = [
    ("Lost sales",
     "Availability issues and execution gaps last longer than they should.\n\n"
     "Problems are visible, but not resolved fast enough to protect revenue."),
    ("Margin leakage",
     "Small pricing, promo, and funding errors compound quietly.\n\n"
     "Each issue looks minor – together, they erode margin week after week."),
    ("Working capital\n& waste",
     "Delayed follow-through leads to excess stock, late orders, and write-offs.\n\n"
     "Capital gets tied up because execution cannot keep pace with reality."),
]

card_w = Inches(3.7)
gap = Inches(0.35)

for i, (title, desc) in enumerate(costs):
    x = Inches(0.8) + int((card_w + gap) * i)
    y = Inches(2.3)
    # Card
    add_rounded_rect(slide, x, y, card_w, Inches(4.2), BG_WHITE, BORDER_CARD)
    # Icon placeholder (teal circle)
    add_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), DUVO_MINT)
    add_circle(slide, x + Inches(0.38), y + Inches(0.38), Inches(0.35), DUVO_TEAL)
    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(0.6),
                 title, FONT_HEADLINE, 20, TEXT_PRIMARY, bold=True)
    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(2.2),
                 desc, FONT_BODY, 12, TEXT_SECONDARY)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5: SOLUTION — "Duvo takes you from insights to getting work done"
# Matches: retail-slide-05.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

add_text_box(slide, Inches(0.6), Inches(0.5), Inches(11), Inches(1.2),
             "Duvo takes you from insights to getting work\ndone in less than two weeks.",
             FONT_HEADLINE, 36, TEXT_PRIMARY, bold=True)

solutions = [
    ("Duvo Clarity", "▶",
     "Captures how work really happens, including steps, exceptions, and "
     "hand-offs, and quantifies the cost of doing nothing.",
     BG_WHITE),
    ("Duvo Case Queue", "▶",
     "Takes each work item from start to finish: gathers inputs, routes "
     "approvals, updates your systems, and keeps an audit trail – until it's done.",
     BG_WHITE),
    ("Finished work", "✓",
     "Every work item ends with a clear outcome: updated systems, "
     "evidence attached, and a proof pack you can forward internally.",
     BG_BLACK),
]

for i, (title, icon, desc, bg) in enumerate(solutions):
    x = Inches(0.8) + int((card_w + gap) * i)
    y = Inches(2.3)
    txt_color = TEXT_WHITE if bg == BG_BLACK else TEXT_PRIMARY
    desc_color = RGBColor(0xCC, 0xCC, 0xCC) if bg == BG_BLACK else TEXT_SECONDARY
    border = None if bg == BG_BLACK else BORDER_CARD

    add_rounded_rect(slide, x, y, card_w, Inches(4.2), bg, border)

    # Icon
    icon_bg = DUVO_YELLOW if bg != BG_BLACK else SUCCESS_GREEN
    add_rounded_rect(slide, x + Inches(0.3), y + Inches(0.3),
                     Inches(0.5), Inches(0.5), icon_bg)
    add_text_box(slide, x + Inches(0.35), y + Inches(0.33),
                 Inches(0.4), Inches(0.4), icon, FONT_BODY, 14, txt_color,
                 alignment=PP_ALIGN.CENTER)

    # Plus / equals between cards
    if i < 2:
        symbol = "+" if i == 0 else "="
        add_text_box(slide, x + card_w + Inches(0.05), y + Inches(1.8),
                     gap, Inches(0.4), symbol, FONT_HEADLINE, 22,
                     TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.6), Inches(0.4),
                 title, FONT_HEADLINE, 18, txt_color, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(2.2),
                 desc, FONT_BODY, 12, desc_color)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6: DUVO CLARITY — Step 1: Capture, Analyze, Quantify
# Matches: retail-slide-06.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

section_badge(slide, Inches(0.6), Inches(0.4), "Duvo Clarity")

add_text_box(slide, Inches(0.6), Inches(1.0), Inches(11), Inches(1.2),
             "Step 1: Duvo Clarity agents capture the work and\n"
             "understand the full process in less than 30 minutes.",
             FONT_HEADLINE, 34, TEXT_PRIMARY, bold=True)

steps = [
    ("Step 1", "Capture",
     "Collect real workflows from operators.\n\n"
     "Video walkthroughs, AI interviews, and existing docs."),
    ("Step 2", "Analyze",
     "Map steps, handoffs, and exceptions.\n\n"
     "Spot bottlenecks, rework loops, and inconsistencies."),
    ("Step 3", "Quantify",
     "Quantify the invisible cost of inefficiencies.\n\n"
     "Attach a € figure to manual work of your team."),
    ("", "Compress\ndiscovery from\nmonths to\nminutes.",
     "Outputs: definition of done\n• baseline + value estimate\n"
     "• automation plan • proof pack template.",
     ),
]

step_w = Inches(2.8)
step_gap = Inches(0.3)

for i, item in enumerate(steps):
    x = Inches(0.8) + int((step_w + step_gap) * i)
    y = Inches(2.8)
    is_dark = (i == 3)
    bg = BG_BLACK if is_dark else BG_WHITE
    border = None if is_dark else BORDER_CARD
    txt = TEXT_WHITE if is_dark else TEXT_PRIMARY
    desc_txt = RGBColor(0xBB, 0xBB, 0xBB) if is_dark else TEXT_SECONDARY

    add_rounded_rect(slide, x, y, step_w, Inches(3.8), bg, border)

    # Step label
    if item[0]:
        add_pill(slide, x + Inches(0.2), y + Inches(0.2),
                 Inches(0.8), Inches(0.25), DUVO_MINT)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.21),
                     Inches(0.7), Inches(0.22), item[0], FONT_BODY, 9,
                     DUVO_TEAL, bold=True)
    else:
        # Checkmark for final card
        add_circle(slide, x + Inches(0.2), y + Inches(0.2),
                   Inches(0.3), DUVO_YELLOW)
        add_text_box(slide, x + Inches(0.23), y + Inches(0.2),
                     Inches(0.3), Inches(0.3), "✓", FONT_BODY, 14,
                     TEXT_PRIMARY, alignment=PP_ALIGN.CENTER)

    # Title
    title_y = y + Inches(0.7) if item[0] else y + Inches(0.7)
    add_text_box(slide, x + Inches(0.2), title_y, step_w - Inches(0.4), Inches(0.8),
                 item[1], FONT_HEADLINE, 18, txt, bold=True)

    # Description
    add_text_box(slide, x + Inches(0.2), title_y + Inches(0.9),
                 step_w - Inches(0.4), Inches(2), item[2], FONT_BODY, 11, desc_txt)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7: CONSULTANT COMPARISON — "Outputs equivalent to top-tier consultants"
# Matches: retail-slide-07.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

section_badge(slide, Inches(0.6), Inches(0.4), "Duvo Clarity")

add_text_box(slide, Inches(0.6), Inches(1.0), Inches(11), Inches(0.8),
             "Outputs equivalent to what top-tier consultants deliver.",
             FONT_HEADLINE, 34, TEXT_PRIMARY, bold=True)

# Three output cards arranged like in the Figma
outputs = [
    ("Complex process diagram with steps",
     Inches(1), Inches(2.2), Inches(5), Inches(3)),
    ("Financial Impact analysis",
     Inches(6.5), Inches(2.2), Inches(5.5), Inches(2)),
    ("Comprehensive transformation proposal",
     Inches(6.5), Inches(4.5), Inches(5.5), Inches(2)),
    ("Actionable process insights",
     Inches(1), Inches(5.5), Inches(5), Inches(1.5)),
]

for label, x, y, w, h in outputs:
    add_rounded_rect(slide, x, y, w, h, BG_WHITE, BORDER_CARD)
    # Icon circle
    add_circle(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.3), BG_SURFACE)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.17),
                 Inches(0.3), Inches(0.25), "📄", FONT_BODY, 11,
                 TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, x + Inches(0.55), y + Inches(0.18), w - Inches(0.7), Inches(0.3),
                 label, FONT_BODY, 13, TEXT_PRIMARY, bold=True)
    # Placeholder lines
    for line_i in range(min(3, int((h - Inches(0.7)) / Inches(0.25)))):
        ly = y + Inches(0.65 + line_i * 0.25)
        add_rect(slide, x + Inches(0.3), ly, w - Inches(0.6), Inches(0.08), BORDER_LIGHT)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8: PRICING COMPARISON — "Big 4 vs Duvo Clarity"
# Matches: retail-slide-08.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

section_badge(slide, Inches(0.6), Inches(0.4), "Duvo Clarity")

add_text_box(slide, Inches(0.6), Inches(1.0), Inches(11), Inches(0.8),
             "From process mapping to transformation proposal\nat a fraction of Big 4 cost.",
             FONT_HEADLINE, 34, TEXT_PRIMARY, bold=True)

# Big 4 card
add_rounded_rect(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4),
                 BG_WHITE, BORDER_CARD)
add_text_box(slide, Inches(1.2), Inches(2.5), Inches(0.4), Inches(0.4),
             "$", FONT_HEADLINE, 24, TEXT_TERTIARY)
add_text_box(slide, Inches(1.2), Inches(3.0), Inches(4.5), Inches(0.5),
             "Big 4 project", FONT_HEADLINE, 22, TEXT_PRIMARY, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.5), Inches(4.5), Inches(0.5),
             "~€400,000", FONT_HEADLINE, 28, PRICE_GREEN, bold=True)
add_text_box(slide, Inches(1.2), Inches(4.2), Inches(4.5), Inches(2),
             "A Big 4 firm would approach this as an 'Operational\n"
             "Excellence/Audit' project delivered through workshops,\n"
             "interviews, and analysis-heavy workstreams.\n\n"
             "Strong documentation and recommendations, but\n"
             "typically slower and harder to scale across multiple\nworkflows.",
             FONT_BODY, 12, TEXT_SECONDARY)

# Duvo card
add_rounded_rect(slide, Inches(7), Inches(2.3), Inches(5.5), Inches(4),
                 BG_WHITE, BORDER_CARD)
# Duvo logo
add_text_box(slide, Inches(7.4), Inches(2.5), Inches(0.4), Inches(0.4),
             "▶", FONT_HEADLINE, 20, DUVO_YELLOW)
add_text_box(slide, Inches(7.4), Inches(3.0), Inches(4.5), Inches(0.5),
             "Duvo Clarity", FONT_HEADLINE, 22, TEXT_PRIMARY, bold=True)
add_text_box(slide, Inches(7.4), Inches(3.5), Inches(4.5), Inches(0.5),
             "TBD", FONT_HEADLINE, 28, PRICE_GREEN, bold=True)
add_text_box(slide, Inches(7.4), Inches(4.2), Inches(4.5), Inches(2),
             'For 1/10th the price you get top-tier "consultant-in-\n'
             'a-box" that turns conversations into complex process\n'
             "diagrams and transformation plans in minutes, not\nmonths.\n\n"
             "Repeatable, run it again for the next process/team.\n\n"
             "(10 runs included free annually; pay as you scale.)",
             FONT_BODY, 12, TEXT_SECONDARY)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9: CASE QUEUE — "Step 2: Duvo Case Queue runs work items"
# Matches: retail-slide-09.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

section_badge(slide, Inches(0.6), Inches(0.4), "Duvo Case Queue")

add_text_box(slide, Inches(0.6), Inches(1.0), Inches(11), Inches(1.2),
             "Step 2: Duvo Case Queue runs work items until\n"
             "everything is done – even for weeks-long processes.",
             FONT_HEADLINE, 34, TEXT_PRIMARY, bold=True)

# Central architecture diagram placeholder
add_rounded_rect(slide, Inches(3.5), Inches(2.8), Inches(6), Inches(3.5),
                 BG_SURFACE, BORDER_LIGHT)
add_text_box(slide, Inches(3.5), Inches(4.0), Inches(6), Inches(1),
             "[Architecture diagram: Duvo Case Queue\nconnecting systems end-to-end]",
             FONT_BODY, 14, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

# Surrounding capability callouts
capabilities = [
    ("Reads your files", "Emails, PDFs, contracts,\nforms, spreadsheets, portal data.",
     Inches(0.3), Inches(2.8)),
    ("Keeps a live queue", "Done + waiting + needs approval — nothing\nslips.",
     Inches(0.3), Inches(4.8)),
    ("Takes action across systems", "SAP, ERPs, supplier portals, shared drives,\nspreadsheets, web applications.",
     Inches(9.8), Inches(2.8)),
    ("Communicates when needed", "Sends emails, prepares reports, and\ncan call suppliers or partners.",
     Inches(9.8), Inches(4.0)),
    ("Delivers finished work (with proof)", "Work is only marked 'done' when the\nresult is written back to your systems.",
     Inches(9.8), Inches(5.2)),
]

for title, desc, x, y in capabilities:
    add_text_box(slide, x, y, Inches(3.2), Inches(0.3),
                 title, FONT_BODY, 12, TEXT_PRIMARY, bold=True)
    add_text_box(slide, x, y + Inches(0.3), Inches(3.2), Inches(0.6),
                 desc, FONT_BODY, 10, TEXT_SECONDARY)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10: BROWSER — "Duvo executes work inside a secure remote browser"
# Matches: retail-slide-10.png
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
content_card_bg(slide)

section_badge(slide, Inches(0.6), Inches(0.4), "Duvo Case Queue")

add_text_box(slide, Inches(0.6), Inches(1.0), Inches(8), Inches(1.5),
             "Duvo executes work\ninside a secure remote\nbrowser",
             FONT_HEADLINE, 36, TEXT_PRIMARY, bold=True)

add_text_box(slide, Inches(0.6), Inches(3.0), Inches(7), Inches(0.6),
             "It operates like your human employee across SAP, spreadsheets,\n"
             "email, supplier portals and internal tools.",
             FONT_BODY, 14, TEXT_SECONDARY)

# Three columns
columns = [
    ("01", "Enterprise\nBrowser",
     "Agents operate through a\nsecure browser session in\nyour systems just like a\nhuman user."),
    ("02", "Start without\nintegrations; add\nAPIs later where\nit's worth it",
     "Agents work directly in\nyour UI.\n\nSkip upfront integrations.\nAdd APIs later where it's\nworth it."),
    ("03", "Enterprise-grade\nisolation and\ngovernance",
     "Ephemeral sandboxes.\nStrict tenant isolation.\nCredential vaults.\nAudit trails."),
]

col_w = Inches(3.5)
for i, (num, title, desc) in enumerate(columns):
    x = Inches(0.8) + int((col_w + Inches(0.4)) * i)
    y = Inches(4.0)
    add_text_box(slide, x, y, Inches(0.5), Inches(0.4),
                 num, FONT_HEADLINE, 28, BORDER_LIGHT, bold=True)
    add_text_box(slide, x, y + Inches(0.5), col_w, Inches(1),
                 title, FONT_HEADLINE, 16, TEXT_PRIMARY, bold=True)
    add_text_box(slide, x, y + Inches(1.6), col_w, Inches(1.5),
                 desc, FONT_BODY, 12, TEXT_SECONDARY)

# Cloud/system logos placeholder (right side)
add_rounded_rect(slide, Inches(9), Inches(1.0), Inches(3.5), Inches(2.5),
                 DUVO_YELLOW_PALE)
add_text_box(slide, Inches(9), Inches(1.8), Inches(3.5), Inches(1),
             "[System logos:\nOffice 365, SAP, Oracle,\nGoogle, Salesforce, etc.]",
             FONT_BODY, 11, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

duvo_footer(slide)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11: CLOSING — Thank you / contact
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
yellow_glow_bg(slide)

# Logo centered large
add_text_box(slide, Inches(3), Inches(2.5), Inches(7), Inches(1),
             "▶  duvo.ai", FONT_HEADLINE, 56, TEXT_PRIMARY, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(3.8), Inches(7), Inches(0.5),
             "Agentic Automation for Retail Operations",
             FONT_HEADLINE, 22, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(5.0), Inches(7), Inches(0.5),
             "www.duvo.ai",
             FONT_BODY, 16, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

duvo_footer(slide)


# ─── Save ───────────────────────────────────────────────────────────────────
output_path = "/root/data/duvo-retail-deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
