#!/usr/bin/env python3
"""
Rekon / wonderful.ai design system deck — v2
Matches the actual rekon-green.vercel.app visual language:
- Layered pastel gradients (pink, purple, blue)
- Warm accent bars (lavender → rose → coral → golden)
- Full-bleed photo placeholders with dark overlays
- DM Sans light (300) headlines, Inter body
- Editorial, photography-heavy feel
- Content-dense slides — no empty section dividers
"""

import sys
sys.path.insert(0, '/root/projects/duvo/deck')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx_helpers import set_slide_bg, add_rect, add_rounded_rect, add_circle, \
    add_text_box, add_multiline_text, add_pill

# ─── Brand tokens ───────────────────────────────────────────────────────────
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE = RGBColor(0xFA, 0xFA, 0xFA)
TEXT_PRIMARY = RGBColor(0x1F, 0x1F, 0x1F)
TEXT_SECONDARY = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_TERTIARY = RGBColor(0x99, 0x99, 0x99)
BORDER = RGBColor(0xE6, 0xE6, 0xE6)

# Dark / inverted
BG_INVERTED = RGBColor(0x0A, 0x0A, 0x0A)
BG_INV_RAISED = RGBColor(0x18, 0x18, 0x18)
BG_INV_ELEVATED = RGBColor(0x23, 0x23, 0x23)
TEXT_INV_PRIMARY = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_INV_SECONDARY = RGBColor(0x88, 0x88, 0x88)
TEXT_INV_TERTIARY = RGBColor(0x55, 0x55, 0x55)
BORDER_INV = RGBColor(0x28, 0x28, 0x28)

# CTA — monochrome, never accent
FILL_ACTION = RGBColor(0x00, 0x00, 0x00)
FILL_ACTION_FG = RGBColor(0xFF, 0xFF, 0xFF)
FILL_ACTION_DARK = RGBColor(0xFF, 0xFF, 0xFF)
FILL_ACTION_DARK_FG = RGBColor(0x00, 0x00, 0x00)

# Accent (info/links ONLY)
ACCENT = RGBColor(0x00, 0x99, 0xFF)

# Semantic
SUCCESS = RGBColor(0x03, 0xA9, 0x7E)
ERROR = RGBColor(0xFF, 0x3F, 0x3F)
WARNING = RGBColor(0xFF, 0xBE, 0x0B)

# ─── GRADIENT palette (from actual rekon site) ─────────────────────────────
PINK = RGBColor(0xF9, 0xA8, 0xD4)
PURPLE = RGBColor(0xC4, 0xB5, 0xFD)
BLUE = RGBColor(0x93, 0xC5, 0xFD)
PINK_LIGHT = RGBColor(0xFD, 0xF2, 0xF8)
PURPLE_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)
BLUE_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)

# Warm accent bar colors (testimonial top borders)
LAVENDER = RGBColor(0xCF, 0xC4, 0xFD)
ROSE = RGBColor(0xF5, 0xC7, 0xCA)
CORAL = RGBColor(0xFF, 0xA7, 0x6A)
GOLDEN = RGBColor(0xFF, 0xDE, 0x97)

# Photo placeholder
PHOTO_WARM = RGBColor(0xE8, 0xE5, 0xE0)
PHOTO_OVERLAY = RGBColor(0x1A, 0x1A, 0x1A)

# Fonts
FONT_DISPLAY = "DM Sans"
FONT_BODY = "Inter"

# Dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ────────────────────────────────────────────────────────────────
# Core helpers imported from pptx_helpers module.


def gradient_band(slide, left, top, total_width, height):
    """Simulate a 4-color warm gradient band (lavender → rose → coral → golden)."""
    w = total_width / 4
    for i, color in enumerate([LAVENDER, ROSE, CORAL, GOLDEN]):
        add_rect(slide, left + int(w * i), top, int(w) + Inches(0.01), height, color)


def photo_placeholder(slide, left, top, width, height, label="Photo"):
    """Warm gray rounded rect simulating an editorial photo."""
    shape = add_rounded_rect(slide, left, top, width, height, PHOTO_WARM)
    add_text_box(slide, left, top + height // 2 - Inches(0.15), width, Inches(0.3),
                 label, FONT_BODY, 11, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)
    return shape


def photo_fullbleed(slide, label="Full-bleed editorial photo"):
    """Full slide warm placeholder with dark overlay on bottom half."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PHOTO_WARM)
    # Dark overlay on lower portion for text
    add_rect(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(4), PHOTO_OVERLAY)
    # Gradient transition strip (semi-dark)
    mid = RGBColor(0x44, 0x44, 0x44)
    add_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.5), mid)


# ─── Build the deck ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

overflow_w = Inches(18)
overflow_left = Inches((13.333 - 18) / 2)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER — Full-bleed photo hero, dark overlay, thin headline
# Matches rekon hero: large photo bg + bottom-aligned DM Sans light text
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
# Photo background
photo_fullbleed(slide)

# Small marker top-left
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "wonderful.", FONT_DISPLAY, 18, WHITE, letter_spacing=-0.8)

# Main headline — bottom-aligned like the rekon hero
add_text_box(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(2),
             "Know what your\ncompetitors charge.", FONT_DISPLAY, 64,
             TEXT_INV_PRIMARY, bold=False, letter_spacing=-3)

# Right-aligned CTA area (like rekon hero right-side)
pill = add_pill(slide, Inches(9.5), Inches(6.3), Inches(2.8), Inches(0.55),
                FILL_ACTION_DARK)
add_text_box(slide, Inches(9.5), Inches(6.35), Inches(2.8), Inches(0.45),
             "Get started", FONT_BODY, 14, FILL_ACTION_DARK_FG,
             alignment=PP_ALIGN.CENTER, letter_spacing=-0.5)

# Bottom metadata
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(3), Inches(0.3),
             "DESIGN SYSTEM · MARCH 2026", FONT_BODY, 10, TEXT_INV_TERTIARY,
             letter_spacing=1.5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2: CAPABILITIES — 3-column feature cards with photos
# Like the rekon features section with cards + imagery
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PAGE_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— CAPABILITIES", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Everything you need to win on pricing.", FONT_DISPLAY, 38,
             TEXT_PRIMARY, bold=False, letter_spacing=-2)

# 3 feature cards with photo placeholders
features = [
    ("Real-time monitoring", "Track competitor prices across 1000+ sources. Get alerts the moment anything changes."),
    ("AI-powered insights", "Machine learning models detect patterns and predict pricing moves before they happen."),
    ("Dynamic repricing", "Automated rules engine adjusts your prices in real-time to stay competitive."),
]

card_w = Inches(3.7)
gap = Inches(0.4)

for i, (title, desc) in enumerate(features):
    x = Inches(0.8) + int((card_w + gap) * i)
    y = Inches(2.4)
    # Card background
    add_rounded_rect(slide, x, y, card_w, Inches(4.5), SURFACE, BORDER)
    # Photo placeholder in card
    photo_placeholder(slide, x + Inches(0.3), y + Inches(0.3),
                      card_w - Inches(0.6), Inches(2.2), "Product screenshot")
    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(2.7), card_w - Inches(0.6), Inches(0.5),
                 title, FONT_DISPLAY, 20, TEXT_PRIMARY, letter_spacing=-1)
    # Description
    add_text_box(slide, x + Inches(0.3), y + Inches(3.3), card_w - Inches(0.6), Inches(1),
                 desc, FONT_BODY, 13, TEXT_SECONDARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3: PLATFORM — Large centered diagram with feature callouts
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PAGE_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— PLATFORM", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Built for scale.", FONT_DISPLAY, 38, TEXT_PRIMARY,
             bold=False, letter_spacing=-2)

# Central diagram placeholder
photo_placeholder(slide, Inches(2.5), Inches(2.3), Inches(8.3), Inches(4.2),
                  "Platform architecture diagram")

# Stat callouts around the edges
stats = [
    ("1000+", "Data sources", Inches(0.3), Inches(2.8)),
    ("< 5 min", "Update cycle", Inches(0.3), Inches(4.0)),
    ("99.9%", "Uptime SLA", Inches(0.3), Inches(5.2)),
    ("50M+", "Products tracked", Inches(11.2), Inches(2.8)),
    ("12", "Markets", Inches(11.2), Inches(4.0)),
    ("3x", "ROI average", Inches(11.2), Inches(5.2)),
]

for value, label, x, y in stats:
    add_text_box(slide, x, y, Inches(1.8), Inches(0.4),
                 value, FONT_DISPLAY, 24, TEXT_PRIMARY, letter_spacing=-1)
    add_text_box(slide, x, y + Inches(0.35), Inches(1.8), Inches(0.25),
                 label, FONT_BODY, 11, TEXT_SECONDARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4: TEAM — Full-bleed photo with overlay text
# Warm office photography feel like the rekon site
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
photo_fullbleed(slide, "Team / office photography")

add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8), Inches(1),
             "Built by pricing experts\nwho get it.", FONT_DISPLAY, 48,
             TEXT_INV_PRIMARY, bold=False, letter_spacing=-2.5)

add_text_box(slide, Inches(0.8), Inches(6.0), Inches(8), Inches(0.6),
             "30+ years of combined experience in retail intelligence,\n"
             "machine learning, and competitive strategy.",
             FONT_BODY, 16, TEXT_INV_SECONDARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5: INDUSTRIES — 3 editorial photo cards
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PAGE_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— INDUSTRIES", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Trusted across verticals.", FONT_DISPLAY, 38, TEXT_PRIMARY,
             bold=False, letter_spacing=-2)

industries = [
    ("Retail & E-commerce", "Thousands of SKUs repriced automatically across all channels."),
    ("Travel & Hospitality", "Dynamic room and ticket pricing that responds to demand in real-time."),
    ("Financial Services", "Fee benchmarking and product pricing intelligence at scale."),
]

for i, (name, desc) in enumerate(industries):
    x = Inches(0.8) + int((card_w + gap) * i)
    y = Inches(2.3)
    # Photo card
    photo_placeholder(slide, x, y, card_w, Inches(2.8), "Industry photo")
    # Gradient accent bar under photo
    bar_colors = [LAVENDER, ROSE, CORAL]
    add_rect(slide, x, y + Inches(2.8), card_w, Inches(0.06), bar_colors[i])
    # Title + desc
    add_text_box(slide, x, y + Inches(3.1), card_w, Inches(0.4),
                 name, FONT_DISPLAY, 20, TEXT_PRIMARY, letter_spacing=-1)
    add_text_box(slide, x, y + Inches(3.6), card_w, Inches(1),
                 desc, FONT_BODY, 13, TEXT_SECONDARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6: FEATURES GRID — 2×2 feature tiles
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, SURFACE)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— FEATURES", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "The details matter.", FONT_DISPLAY, 38, TEXT_PRIMARY,
             bold=False, letter_spacing=-2)

features_grid = [
    ("Automated collection", "Scrape, API, and feed ingestion from 1000+ sources with built-in deduplication."),
    ("Smart matching", "ML-powered product matching across retailers — handles variants, bundles, and regional SKUs."),
    ("Alerting engine", "Slack, email, and webhook alerts with configurable thresholds per product, category, or competitor."),
    ("Custom dashboards", "Build views that matter to your team. Drag-and-drop with real-time data refresh."),
]

tile_w = Inches(5.8)
tile_h = Inches(2.1)

for i, (title, desc) in enumerate(features_grid):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + int(col * (tile_w + Inches(0.4)))
    y = Inches(2.3) + int(row * (tile_h + Inches(0.4)))
    # Tile
    add_rounded_rect(slide, x, y, tile_w, tile_h, WHITE, BORDER)
    # Colored dot
    dot_colors = [LAVENDER, ROSE, CORAL, GOLDEN]
    add_circle(slide, x + Inches(0.4), y + Inches(0.4), Inches(0.15), dot_colors[i])
    # Title
    add_text_box(slide, x + Inches(0.75), y + Inches(0.3), tile_w - Inches(1), Inches(0.4),
                 title, FONT_DISPLAY, 20, TEXT_PRIMARY, letter_spacing=-1)
    # Description
    add_text_box(slide, x + Inches(0.4), y + Inches(0.9), tile_w - Inches(0.8), Inches(1),
                 desc, FONT_BODY, 13, TEXT_SECONDARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7: TESTIMONIALS — Dark bg, gradient accent bars on white cards
# Matches rekon testimonial section exactly
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, BG_INVERTED)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— TESTIMONIALS", FONT_BODY, 11, TEXT_INV_TERTIARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "What our clients say.", FONT_DISPLAY, 38, TEXT_INV_PRIMARY,
             bold=False, letter_spacing=-2)

# Full-width gradient band
gradient_band(slide, Inches(0.8), Inches(2.15), Inches(11.7), Inches(0.06))

testimonials = [
    ('"wonderful. cut our price analysis time by 80%."', "— Sarah Chen, VP Pricing, RetailCo"),
    ('"The accuracy of their matching engine is unreal."', "— Marcus Webb, Head of Strategy, TravelMax"),
    ('"Finally a tool that scales with our catalog."', "— Anika Patel, Director of E-commerce, FinServ"),
]

quote_card_w = Inches(3.7)
for i, (quote, author) in enumerate(testimonials):
    x = Inches(0.8) + int((quote_card_w + Inches(0.3)) * i)
    y = Inches(2.6)
    # White card
    add_rounded_rect(slide, x, y, quote_card_w, Inches(3.5), WHITE)
    # Colored top accent bar
    bar_colors = [LAVENDER, ROSE, CORAL]
    add_rect(slide, x + Inches(0.3), y + Inches(0.3), Inches(2), Inches(0.05),
             bar_colors[i])
    # Quote
    add_text_box(slide, x + Inches(0.3), y + Inches(0.65), quote_card_w - Inches(0.6),
                 Inches(1.8), quote, FONT_DISPLAY, 18, TEXT_PRIMARY, letter_spacing=-0.8)
    # Author
    add_text_box(slide, x + Inches(0.3), y + Inches(2.7), quote_card_w - Inches(0.6),
                 Inches(0.4), author, FONT_BODY, 12, TEXT_SECONDARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8: DESIGN SYSTEM — Typography + color overview
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PAGE_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— DESIGN SYSTEM", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Typography & color", FONT_DISPLAY, 38, TEXT_PRIMARY,
             bold=False, letter_spacing=-2)

# Left: Typography samples
add_text_box(slide, Inches(0.8), Inches(2.2), Inches(3), Inches(0.3),
             "DM Sans · light (300)", FONT_BODY, 10, ACCENT, letter_spacing=1)

type_samples = [
    ("Display large", 44, -2.5),
    ("Heading one", 28, -1.5),
    ("Heading two", 22, -1),
    ("Heading three", 18, -0.8),
]
y_pos = 2.6
for text, size, ls in type_samples:
    add_text_box(slide, Inches(0.8), Inches(y_pos), Inches(5), Inches(0.7),
                 text, FONT_DISPLAY, size, TEXT_PRIMARY, bold=False, letter_spacing=ls)
    add_text_box(slide, Inches(0.8), Inches(y_pos + 0.5), Inches(5), Inches(0.2),
                 f"{size}px / weight 300 / {ls}px", FONT_BODY, 9, TEXT_TERTIARY)
    y_pos += 0.75

add_text_box(slide, Inches(0.8), Inches(y_pos + 0.2), Inches(5), Inches(0.3),
             "Inter · regular (400)", FONT_BODY, 10, ACCENT, letter_spacing=1)
add_text_box(slide, Inches(0.8), Inches(y_pos + 0.5), Inches(5), Inches(0.3),
             "Body text at 17px for paragraphs and content", FONT_BODY, 17, TEXT_PRIMARY)
add_text_box(slide, Inches(0.8), Inches(y_pos + 0.85), Inches(5), Inches(0.3),
             "Caption at 14px for metadata and labels", FONT_BODY, 14, TEXT_SECONDARY)

# Right: Color palette
add_text_box(slide, Inches(7), Inches(2.2), Inches(3), Inches(0.3),
             "Core palette", FONT_BODY, 10, ACCENT, letter_spacing=1)

mono = [
    (RGBColor(0x00, 0x00, 0x00), "#000", "Action"),
    (RGBColor(0x0A, 0x0A, 0x0A), "#0A0", "Inverted"),
    (RGBColor(0x1F, 0x1F, 0x1F), "#1F1", "Primary"),
    (RGBColor(0x6B, 0x6B, 0x6B), "#6B6", "Secondary"),
    (RGBColor(0xE6, 0xE6, 0xE6), "#E6E", "Border"),
    (RGBColor(0xFF, 0xFF, 0xFF), "#FFF", "Page"),
]

for i, (c, h, label) in enumerate(mono):
    x = Inches(7 + i * 1)
    add_rounded_rect(slide, x, Inches(2.65), Inches(0.7), Inches(0.7), c,
                     BORDER if h in ("#E6E", "#FFF") else None)
    add_text_box(slide, x, Inches(3.45), Inches(0.7), Inches(0.2),
                 label, FONT_BODY, 8, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)

# Gradient palette
add_text_box(slide, Inches(7), Inches(4.0), Inches(3), Inches(0.3),
             "Gradient palette", FONT_BODY, 10, ACCENT, letter_spacing=1)

grad_colors = [
    (PINK, "Pink"), (PURPLE, "Purple"), (BLUE, "Blue"),
    (LAVENDER, "Lavender"), (ROSE, "Rose"), (CORAL, "Coral"), (GOLDEN, "Golden"),
]
for i, (c, name) in enumerate(grad_colors):
    col = i % 4
    row = i // 4
    x = Inches(7 + col * 1.55)
    y = Inches(4.4 + row * 1.1)
    add_rounded_rect(slide, x, y, Inches(1.3), Inches(0.6), c)
    add_text_box(slide, x, y + Inches(0.65), Inches(1.3), Inches(0.2),
                 name, FONT_BODY, 9, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# Key rule at bottom
add_rect(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.007), BORDER)
add_text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.3),
             "DM Sans 300 is the signature. Never bold headlines. Blue accent is for info/focus — CTAs are black/white.",
             FONT_BODY, 11, TEXT_TERTIARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9: COMPONENTS — Buttons, cards, inputs, badges
# Dense showcase of the 194-component system
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PAGE_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— COMPONENTS", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "194 components", FONT_DISPLAY, 38, TEXT_PRIMARY,
             bold=False, letter_spacing=-2)

# Row 1: Buttons
add_text_box(slide, Inches(0.8), Inches(2.1), Inches(2), Inches(0.3),
             "Buttons", FONT_BODY, 11, TEXT_TERTIARY, letter_spacing=1)

# Primary pill
add_pill(slide, Inches(0.8), Inches(2.5), Inches(2.3), Inches(0.5), FILL_ACTION)
add_text_box(slide, Inches(0.8), Inches(2.53), Inches(2.3), Inches(0.44),
             "Get started", FONT_BODY, 14, FILL_ACTION_FG,
             alignment=PP_ALIGN.CENTER, letter_spacing=-0.5)

# Secondary pill
add_pill(slide, Inches(3.4), Inches(2.5), Inches(2.3), Inches(0.5), WHITE, BORDER)
add_text_box(slide, Inches(3.4), Inches(2.53), Inches(2.3), Inches(0.44),
             "Learn more", FONT_BODY, 14, TEXT_PRIMARY,
             alignment=PP_ALIGN.CENTER, letter_spacing=-0.5)

# Ghost
add_text_box(slide, Inches(6), Inches(2.53), Inches(2), Inches(0.44),
             "Ghost link →", FONT_BODY, 14, TEXT_PRIMARY, letter_spacing=-0.5)

# Danger pill
add_pill(slide, Inches(8.2), Inches(2.5), Inches(1.8), Inches(0.5), ERROR)
add_text_box(slide, Inches(8.2), Inches(2.53), Inches(1.8), Inches(0.44),
             "Delete", FONT_BODY, 14, WHITE, alignment=PP_ALIGN.CENTER)

# Dark context buttons
add_rect(slide, Inches(10.3), Inches(2.3), Inches(2.7), Inches(0.9), BG_INVERTED)
add_pill(slide, Inches(10.5), Inches(2.5), Inches(2.3), Inches(0.5), FILL_ACTION_DARK)
add_text_box(slide, Inches(10.5), Inches(2.53), Inches(2.3), Inches(0.44),
             "Get started", FONT_BODY, 14, FILL_ACTION_DARK_FG,
             alignment=PP_ALIGN.CENTER, letter_spacing=-0.5)

# Row 2: Cards
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(2), Inches(0.3),
             "Cards", FONT_BODY, 11, TEXT_TERTIARY, letter_spacing=1)

# Example card
add_rounded_rect(slide, Inches(0.8), Inches(3.7), Inches(3.5), Inches(2.3), WHITE, BORDER)
photo_placeholder(slide, Inches(1.0), Inches(3.9), Inches(3.1), Inches(1.1), "Image")
add_text_box(slide, Inches(1.0), Inches(5.15), Inches(3.1), Inches(0.3),
             "Card title", FONT_DISPLAY, 16, TEXT_PRIMARY, letter_spacing=-0.5)
add_text_box(slide, Inches(1.0), Inches(5.45), Inches(3.1), Inches(0.3),
             "Description text goes here.", FONT_BODY, 12, TEXT_SECONDARY)

# Stat card
add_rounded_rect(slide, Inches(4.7), Inches(3.7), Inches(2.5), Inches(2.3), WHITE, BORDER)
add_text_box(slide, Inches(5.0), Inches(3.9), Inches(2), Inches(0.3),
             "Products tracked", FONT_BODY, 11, TEXT_SECONDARY)
add_text_box(slide, Inches(5.0), Inches(4.35), Inches(2), Inches(0.5),
             "50M+", FONT_DISPLAY, 36, TEXT_PRIMARY, letter_spacing=-1.5)
add_circle(slide, Inches(5.0), Inches(5.1), Inches(0.12), SUCCESS)
add_text_box(slide, Inches(5.2), Inches(5.05), Inches(2), Inches(0.3),
             "+12.3% this month", FONT_BODY, 12, SUCCESS)

# Row 2b: Inputs + badges
add_text_box(slide, Inches(7.6), Inches(3.3), Inches(2), Inches(0.3),
             "Inputs & badges", FONT_BODY, 11, TEXT_TERTIARY, letter_spacing=1)

# Input field
add_rounded_rect(slide, Inches(7.6), Inches(3.7), Inches(3), Inches(0.5), WHITE, BORDER)
add_text_box(slide, Inches(7.8), Inches(3.73), Inches(2.8), Inches(0.44),
             "Search products...", FONT_BODY, 13, TEXT_TERTIARY)

# Badges
badge_data = [
    ("ACTIVE", SUCCESS), ("WARNING", WARNING), ("ERROR", ERROR), ("INFO", ACCENT),
]
for i, (label, dot_color) in enumerate(badge_data):
    x = Inches(7.6 + i * 1.5)
    y = Inches(4.5)
    add_pill(slide, x, y, Inches(1.3), Inches(0.35), SURFACE)
    add_circle(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.12), dot_color)
    add_text_box(slide, x + Inches(0.35), y + Inches(0.03), Inches(0.85), Inches(0.3),
                 label, FONT_BODY, 9, TEXT_PRIMARY, letter_spacing=0.8)

# Toggle + checkbox sketches
add_rounded_rect(slide, Inches(7.6), Inches(5.1), Inches(0.5), Inches(0.28), FILL_ACTION)
add_circle(slide, Inches(7.88), Inches(5.12), Inches(0.22), WHITE)
add_text_box(slide, Inches(8.3), Inches(5.1), Inches(2), Inches(0.28),
             "Toggle on", FONT_BODY, 12, TEXT_PRIMARY)

add_rounded_rect(slide, Inches(7.6), Inches(5.6), Inches(0.5), Inches(0.28),
                 RGBColor(0x33, 0x33, 0x33))
add_circle(slide, Inches(7.62), Inches(5.62), Inches(0.22), TEXT_SECONDARY)
add_text_box(slide, Inches(8.3), Inches(5.6), Inches(2), Inches(0.28),
             "Toggle off", FONT_BODY, 12, TEXT_SECONDARY)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.3),
             "72 atoms · 85 molecules · 7 organisms · 18 page blocks · 8 templates · 4 layout components",
             FONT_BODY, 11, TEXT_TERTIARY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10: GRADIENT CTA — Layered pastel gradient background
# Matches the rekon CTA section with pink-purple-blue wash
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PINK_LIGHT)

# Layered pastel blocks to simulate the radial gradient
add_rect(slide, Inches(0), Inches(0), Inches(5), SLIDE_H, PINK_LIGHT)
add_rect(slide, Inches(4.5), Inches(0), Inches(5), SLIDE_H, BLUE_LIGHT)
add_rect(slide, Inches(9), Inches(0), Inches(4.5), SLIDE_H, PURPLE_LIGHT)
# Soft overlap blends
add_rect(slide, Inches(3.5), Inches(1), Inches(3), Inches(5),
         RGBColor(0xF6, 0xF0, 0xFE))  # pink-blue transition
add_rect(slide, Inches(7.5), Inches(0.5), Inches(3), Inches(6),
         RGBColor(0xEE, 0xF2, 0xFF))  # blue-purple transition

# Centered content
add_text_box(slide, Inches(2), Inches(2.0), Inches(9.3), Inches(2),
             "Ready to see what\nyour competitors charge?", FONT_DISPLAY, 52,
             TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.CENTER, letter_spacing=-2.5)

add_text_box(slide, Inches(3), Inches(4.3), Inches(7.3), Inches(0.6),
             "Start your free trial today. No credit card required.",
             FONT_BODY, 18, TEXT_SECONDARY, alignment=PP_ALIGN.CENTER)

# CTA buttons centered
btn_x = Inches(4.2)
add_pill(slide, btn_x, Inches(5.3), Inches(2.8), Inches(0.6), FILL_ACTION)
add_text_box(slide, btn_x, Inches(5.35), Inches(2.8), Inches(0.5),
             "Start free trial", FONT_BODY, 16, FILL_ACTION_FG,
             alignment=PP_ALIGN.CENTER, letter_spacing=-0.5)

add_pill(slide, btn_x + Inches(3.2), Inches(5.3), Inches(2.8), Inches(0.6),
         WHITE, BORDER)
add_text_box(slide, btn_x + Inches(3.2), Inches(5.35), Inches(2.8), Inches(0.5),
             "Book a demo", FONT_BODY, 16, TEXT_PRIMARY,
             alignment=PP_ALIGN.CENTER, letter_spacing=-0.5)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11: FAQ — Clean accordion-style layout
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, PAGE_BG)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3),
             "— FAQ", FONT_BODY, 11, TEXT_SECONDARY, letter_spacing=1.5)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.8),
             "Common questions.", FONT_DISPLAY, 38, TEXT_PRIMARY,
             bold=False, letter_spacing=-2)

faqs = [
    ("How does the pricing data collection work?",
     "We combine web scraping, API integrations, and data feed partnerships to monitor over 1000 sources in near real-time."),
    ("What accuracy level can I expect?",
     "Our ML matching engine achieves 97%+ accuracy on product identification, with human QA on edge cases."),
    ("How quickly does pricing data update?",
     "Standard plans refresh every 15 minutes. Enterprise gets sub-5-minute updates with webhook notifications."),
    ("Can I integrate with my existing tools?",
     "Yes — we offer REST APIs, Zapier connectors, and native integrations with Shopify, BigCommerce, and SAP."),
    ("What does onboarding look like?",
     "A dedicated success manager maps your competitive landscape in week one. Most teams are live within 5 days."),
]

for i, (question, answer) in enumerate(faqs):
    y = Inches(2.2 + i * 1.0)
    # Divider
    add_rect(slide, Inches(0.8), y - Inches(0.05), Inches(11.7), Inches(0.007), BORDER)
    # Question
    add_text_box(slide, Inches(0.8), y + Inches(0.05), Inches(5.5), Inches(0.4),
                 question, FONT_DISPLAY, 17, TEXT_PRIMARY, letter_spacing=-0.6)
    # Answer
    add_text_box(slide, Inches(6.8), y + Inches(0.05), Inches(5.7), Inches(0.6),
                 answer, FONT_BODY, 13, TEXT_SECONDARY)
    # Expand icon placeholder
    add_text_box(slide, Inches(12.3), y + Inches(0.05), Inches(0.5), Inches(0.3),
                 "+", FONT_DISPLAY, 18, TEXT_TERTIARY, alignment=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12: CLOSING — Dark footer with gradient band and oversized text
# ═════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank)
set_slide_bg(slide, BG_INVERTED)

# Gradient band near top
gradient_band(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.08))

# Oversized "wonderful." — the signature
add_text_box(slide, overflow_left, Inches(2.0), overflow_w, Inches(3),
             "wonderful.", FONT_DISPLAY, 120, TEXT_INV_PRIMARY, bold=False,
             alignment=PP_ALIGN.CENTER, letter_spacing=-4)

# Tagline
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "Know what your competitors charge.", FONT_DISPLAY, 22,
             TEXT_INV_SECONDARY, alignment=PP_ALIGN.CENTER, letter_spacing=-1)

# File references
add_text_box(slide, Inches(0.8), Inches(6.0), Inches(8), Inches(0.3),
             "tokens.css · typography.md · spacing.md · components/",
             FONT_BODY, 11, TEXT_INV_TERTIARY)

# Bottom gradient accent line
gradient_band(slide, Inches(0), Inches(6.8), SLIDE_W, Inches(0.06))

# Bottom-right
add_text_box(slide, Inches(7), Inches(6.2), Inches(5.5), Inches(0.3),
             "wonderful.ai", FONT_DISPLAY, 18, TEXT_INV_PRIMARY,
             alignment=PP_ALIGN.RIGHT, letter_spacing=-0.8)
add_text_box(slide, Inches(7), Inches(6.5), Inches(5.5), Inches(0.3),
             "194 components · 68 tokens · Light & dark",
             FONT_BODY, 11, TEXT_INV_SECONDARY, alignment=PP_ALIGN.RIGHT)


# ─── Save ───────────────────────────────────────────────────────────────────
output_path = "/root/data/rekon-design-system-deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
